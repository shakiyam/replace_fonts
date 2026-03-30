import tempfile
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation

from font_policy import FontPolicy, load_font_policy, update_theme_fonts
from replace_fonts import log, main, process_pptx_file

POLICY_PATH = Path(__file__).parent / "policy.yaml"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
EA_SCRIPTS = ("Jpan", "Hang", "Hans", "Hant")

EXPECTED_POLICY = FontPolicy(
    major_latin="Arial", major_ea="メイリオ",
    minor_latin="Arial", minor_ea="メイリオ",
)
ORIGINAL_THEME = FontPolicy(
    major_latin="Calibri Light", major_ea="",
    minor_latin="Constantia", minor_ea="",
)
SAMPLE_PPTX = "sample1.pptx"


def _get_font_scheme(pptx_path: Path) -> etree._Element:
    prs = Presentation(str(pptx_path))
    for part in prs.part.package.iter_parts():
        if "theme" in part.content_type:
            root = etree.fromstring(part.blob)
            scheme = root.find(f".//{{{A_NS}}}fontScheme")
            if scheme is not None:
                return scheme
    msg = "No fontScheme found"
    raise ValueError(msg)


def test_load_valid() -> None:
    policy = load_font_policy(POLICY_PATH)
    assert policy == EXPECTED_POLICY


def test_load_empty_yaml() -> None:
    with tempfile.TemporaryDirectory() as tmpdir:
        empty = Path(tmpdir) / "empty.yaml"
        empty.write_text("")
        with pytest.raises(ValueError, match="mapping"):
            load_font_policy(empty)


def test_load_missing_key() -> None:
    with tempfile.TemporaryDirectory() as tmpdir:
        bad = Path(tmpdir) / "bad.yaml"
        bad.write_text("theme_fonts:\n  major:\n    latin: X\n")
        with pytest.raises(ValueError, match="missing"):
            load_font_policy(bad)


def test_load_file_not_found() -> None:
    with pytest.raises(FileNotFoundError):
        load_font_policy(Path("/nonexistent/policy.yaml"))


def test_update_theme_fonts(workspace: tuple[Path, Path]) -> None:
    work_dir, _ = workspace
    pptx_path = work_dir / SAMPLE_PPTX

    scheme = _get_font_scheme(pptx_path)
    major = scheme.find(f"{{{A_NS}}}majorFont")
    minor = scheme.find(f"{{{A_NS}}}minorFont")
    assert major.find(f"{{{A_NS}}}latin").get("typeface") == ORIGINAL_THEME.major_latin
    assert major.find(f"{{{A_NS}}}ea").get("typeface") == ORIGINAL_THEME.major_ea
    assert minor.find(f"{{{A_NS}}}latin").get("typeface") == ORIGINAL_THEME.minor_latin
    assert minor.find(f"{{{A_NS}}}ea").get("typeface") == ORIGINAL_THEME.minor_ea

    prs = Presentation(str(pptx_path))
    log_path = pptx_path.with_suffix(".log")
    with open(log_path, "w") as log_file:
        update_theme_fonts(prs, EXPECTED_POLICY, log_file, log)
    prs.save(str(pptx_path))

    scheme = _get_font_scheme(pptx_path)
    major = scheme.find(f"{{{A_NS}}}majorFont")
    minor = scheme.find(f"{{{A_NS}}}minorFont")
    assert major.find(f"{{{A_NS}}}latin").get("typeface") == EXPECTED_POLICY.major_latin
    assert major.find(f"{{{A_NS}}}ea").get("typeface") == EXPECTED_POLICY.major_ea
    assert minor.find(f"{{{A_NS}}}latin").get("typeface") == EXPECTED_POLICY.minor_latin
    assert minor.find(f"{{{A_NS}}}ea").get("typeface") == EXPECTED_POLICY.minor_ea

    for script in EA_SCRIPTS:
        major_font = major.find(f"{{{A_NS}}}font[@script='{script}']")
        assert major_font.get("typeface") == EXPECTED_POLICY.major_ea
        minor_font = minor.find(f"{{{A_NS}}}font[@script='{script}']")
        assert minor_font.get("typeface") == EXPECTED_POLICY.minor_ea


def test_update_theme_fonts_logs(workspace: tuple[Path, Path]) -> None:
    work_dir, _ = workspace
    pptx_path = work_dir / SAMPLE_PPTX

    prs = Presentation(str(pptx_path))
    log_path = pptx_path.with_suffix(".log")
    with open(log_path, "w") as log_file:
        update_theme_fonts(prs, EXPECTED_POLICY, log_file, log)

    log_content = log_path.read_text()
    old = ORIGINAL_THEME.major_latin
    assert f'Update theme major latin from "{old}"' in log_content
    for script in EA_SCRIPTS:
        assert f"ea script {script}" in log_content


def test_update_theme_fonts_no_change(workspace: tuple[Path, Path]) -> None:
    work_dir, _ = workspace
    pptx_path = work_dir / SAMPLE_PPTX

    prs = Presentation(str(pptx_path))
    log_path = pptx_path.with_suffix(".log")
    with open(log_path, "w") as log_file:
        update_theme_fonts(prs, EXPECTED_POLICY, log_file, log)

    with open(log_path, "w") as log_file:
        update_theme_fonts(prs, EXPECTED_POLICY, log_file, log)

    log_content = log_path.read_text()
    assert "Update theme" not in log_content


def test_dry_run_does_not_modify(workspace: tuple[Path, Path]) -> None:
    work_dir, _ = workspace
    pptx_path = work_dir / SAMPLE_PPTX
    original_content = pptx_path.read_bytes()

    process_pptx_file(
        pptx_path,
        preserve_code_fonts=True,
        dry_run=True,
        font_policy_path=POLICY_PATH,
    )

    assert pptx_path.read_bytes() == original_content


def test_normal_run_updates_theme(workspace: tuple[Path, Path]) -> None:
    work_dir, _ = workspace
    pptx_path = work_dir / SAMPLE_PPTX

    process_pptx_file(
        pptx_path,
        preserve_code_fonts=True,
        font_policy_path=POLICY_PATH,
    )

    scheme = _get_font_scheme(pptx_path)
    major = scheme.find(f"{{{A_NS}}}majorFont")
    assert major.find(f"{{{A_NS}}}latin").get("typeface") == EXPECTED_POLICY.major_latin


def test_with_code(workspace: tuple[Path, Path]) -> None:
    work_dir, _ = workspace
    pptx_path = work_dir / SAMPLE_PPTX

    process_pptx_file(
        pptx_path,
        preserve_code_fonts=True,
        font_policy_path=POLICY_PATH,
    )

    log_content = pptx_path.with_suffix(".log").read_text()
    assert "Update theme" in log_content
    assert "Preserve" in log_content or "Replace" in log_content


def test_cli(
    workspace: tuple[Path, Path], monkeypatch: pytest.MonkeyPatch
) -> None:
    work_dir, _ = workspace
    pptx_path = work_dir / SAMPLE_PPTX

    args = [
        "replace_fonts.py",
        "--font-policy",
        str(POLICY_PATH),
        str(pptx_path),
    ]
    monkeypatch.setattr("sys.argv", args)

    exit_code = main()
    assert exit_code == 0

    scheme = _get_font_scheme(pptx_path)
    major = scheme.find(f"{{{A_NS}}}majorFont")
    assert major.find(f"{{{A_NS}}}latin").get("typeface") == EXPECTED_POLICY.major_latin


def test_cli_dry_run(
    workspace: tuple[Path, Path], monkeypatch: pytest.MonkeyPatch
) -> None:
    work_dir, _ = workspace
    pptx_path = work_dir / SAMPLE_PPTX
    original_content = pptx_path.read_bytes()

    args = [
        "replace_fonts.py",
        "--font-policy",
        str(POLICY_PATH),
        "--dry-run",
        str(pptx_path),
    ]
    monkeypatch.setattr("sys.argv", args)

    exit_code = main()
    assert exit_code == 0
    assert pptx_path.read_bytes() == original_content

    log_content = pptx_path.with_suffix(".log").read_text()
    assert "Update theme" in log_content
