from enum import Enum

from lxml.etree import _Element
from pptx.oxml import CT_TextCharacterProperties  # type: ignore[attr-defined]
from pptx.oxml.ns import qn
from pptx.presentation import Presentation as PresentationType
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.slide import SlideMasters, Slides
from pptx.text.text import TextFrame

from logger import Logger


class ThemeFont(Enum):
    MAJOR = "major"
    MINOR = "minor"


class FontScript(Enum):
    LATIN = "latin"
    EAST_ASIAN = "east asian"


FONT_MAPPINGS = {
    FontScript.LATIN: {
        ThemeFont.MAJOR: "+mj-lt",
        ThemeFont.MINOR: "+mn-lt",
    },
    FontScript.EAST_ASIAN: {
        ThemeFont.MAJOR: "+mj-ea",
        ThemeFont.MINOR: "+mn-ea",
    },
}

FONT_ELEMENT_MAPPINGS = [
    (qn("a:latin"), FontScript.LATIN),
    (qn("a:ea"), FontScript.EAST_ASIAN),
]

PRESERVED_CODE_FONT = "Consolas"
CODE_FONTS_TO_REPLACE = ("Courier New",)


def log_font_action(
    theme_font: ThemeFont,
    font_script: FontScript,
    current_font: str,
    new_font: str | None,
    logger: Logger,
    element_text: str | None = None,
) -> None:
    if new_font:
        message = (
            f"Replace {theme_font.value} {font_script.value} "
            f"from {current_font} to {new_font}"
        )
    else:
        message = f"Preserve {theme_font.value} {font_script.value} as {current_font}"
    logger.log(message, element_text)


def replace_font_element(
    element: _Element,
    theme_font: ThemeFont,
    font_script: FontScript,
    preserve_code_fonts: bool,
    logger: Logger,
    element_text: str | None = None,
) -> None:
    default_font = FONT_MAPPINGS[font_script][theme_font]
    current_font = element.get("typeface")
    if preserve_code_fonts and current_font == PRESERVED_CODE_FONT:
        log_font_action(
            theme_font, font_script, current_font, None, logger, element_text
        )
    elif preserve_code_fonts and current_font in CODE_FONTS_TO_REPLACE:
        element.set("typeface", PRESERVED_CODE_FONT)
        log_font_action(
            theme_font,
            font_script,
            current_font,
            PRESERVED_CODE_FONT,
            logger,
            element_text,
        )
    elif current_font != default_font:
        element.set("typeface", default_font)
        log_font_action(
            theme_font, font_script, current_font, default_font, logger, element_text
        )


def replace_properties_fonts(
    properties: CT_TextCharacterProperties,
    theme_font: ThemeFont,
    preserve_code_fonts: bool,
    logger: Logger,
    element_text: str | None = None,
) -> None:
    for qname, font_script in FONT_ELEMENT_MAPPINGS:
        element = properties.find(qname)
        if element is not None:
            replace_font_element(
                element,
                theme_font,
                font_script,
                preserve_code_fonts,
                logger,
                element_text,
            )


def replace_text_frame_fonts(
    text_frame: TextFrame,
    theme_font: ThemeFont,
    preserve_code_fonts: bool,
    logger: Logger,
) -> None:
    lst_style = text_frame._element.find(qn("a:lstStyle"))
    if lst_style is not None:
        for level_ppr in lst_style:
            def_rpr = level_ppr.find(qn("a:defRPr"))
            if def_rpr is not None:
                replace_properties_fonts(
                    def_rpr, theme_font, preserve_code_fonts, logger
                )
    for paragraph in text_frame.paragraphs:
        if (
            paragraph._element.pPr is not None
            and paragraph._element.pPr.defRPr is not None
        ):
            replace_properties_fonts(
                paragraph._element.pPr.defRPr,
                theme_font,
                preserve_code_fonts,
                logger,
            )
        for run in paragraph.runs:
            run_text = run.text.strip()
            replace_properties_fonts(
                run.font._element,
                theme_font,
                preserve_code_fonts,
                logger,
                run_text,
            )
        for br in paragraph._element.findall(qn("a:br")):
            br_rpr = br.find(qn("a:rPr"))
            if br_rpr is not None:
                replace_properties_fonts(
                    br_rpr, theme_font, preserve_code_fonts, logger
                )
        if paragraph._element.endParaRPr is not None:
            replace_properties_fonts(
                paragraph._element.endParaRPr,
                theme_font,
                preserve_code_fonts,
                logger,
            )


def replace_shape_text_fonts(
    shape: Shape, preserve_code_fonts: bool, logger: Logger
) -> None:
    placeholder = shape.element.find(f".//{qn('p:ph')}")
    if placeholder is not None and placeholder.get("type") in ["ctrTitle", "title"]:
        theme_font = ThemeFont.MAJOR
    else:
        theme_font = ThemeFont.MINOR
    replace_text_frame_fonts(
        shape.text_frame, theme_font, preserve_code_fonts, logger
    )


def replace_table_fonts(
    shape: GraphicFrame, preserve_code_fonts: bool, logger: Logger
) -> None:
    for row in shape.table.rows:
        for cell in row.cells:
            replace_text_frame_fonts(
                cell.text_frame, ThemeFont.MINOR, preserve_code_fonts, logger
            )


def replace_chart_fonts(
    shape: GraphicFrame, preserve_code_fonts: bool, logger: Logger
) -> None:
    for qname, font_script in FONT_ELEMENT_MAPPINGS:
        for element in shape.chart.element.findall(f".//{qname}"):
            replace_font_element(
                element, ThemeFont.MINOR, font_script, preserve_code_fonts, logger
            )


def replace_group_fonts(
    shape: GroupShape, preserve_code_fonts: bool, logger: Logger
) -> None:
    for item in shape.shapes:
        replace_shape_fonts(item, preserve_code_fonts, logger)


def replace_generic_graphicframe_fonts(
    shape: GraphicFrame, preserve_code_fonts: bool, logger: Logger
) -> None:
    for qname, font_script in FONT_ELEMENT_MAPPINGS:
        for element in shape.element.findall(f".//{qname}"):
            replace_font_element(
                element, ThemeFont.MINOR, font_script, preserve_code_fonts, logger
            )


def replace_shape_fonts(
    shape: BaseShape, preserve_code_fonts: bool, logger: Logger
) -> None:
    if isinstance(shape, Shape):
        replace_shape_text_fonts(shape, preserve_code_fonts, logger)
    elif isinstance(shape, GraphicFrame) and shape.has_table:
        replace_table_fonts(shape, preserve_code_fonts, logger)
    elif isinstance(shape, GraphicFrame) and shape.has_chart:
        replace_chart_fonts(shape, preserve_code_fonts, logger)
    elif isinstance(shape, GraphicFrame):
        replace_generic_graphicframe_fonts(shape, preserve_code_fonts, logger)
    elif isinstance(shape, GroupShape):
        replace_group_fonts(shape, preserve_code_fonts, logger)


def process_slides(slides: Slides, preserve_code_fonts: bool, logger: Logger) -> None:
    for i, slide in enumerate(slides):
        logger.log(f"--- Slide {i + 1} ---")
        for shape in slide.shapes:
            replace_shape_fonts(shape, preserve_code_fonts, logger)
        if slide.has_notes_slide:
            logger.log(f"--- Notes Slide {i + 1} ---")
            for shape in slide.notes_slide.shapes:
                replace_shape_fonts(shape, preserve_code_fonts, logger)


def process_slide_masters(
    slide_masters: SlideMasters, preserve_code_fonts: bool, logger: Logger
) -> None:
    for i, slide_master in enumerate(slide_masters):
        logger.log(f"--- Slide Master {i + 1} ---")
        text_styles = slide_master.element.find(qn("p:txStyles"))
        if text_styles is not None:
            for text_style in text_styles:
                if text_style.tag == qn("p:titleStyle"):
                    theme_font = ThemeFont.MAJOR
                else:
                    theme_font = ThemeFont.MINOR
                for list_style in text_style:
                    if isinstance(list_style, CT_TextCharacterProperties):
                        replace_properties_fonts(
                            list_style,
                            theme_font,
                            preserve_code_fonts,
                            logger,
                        )
                    else:
                        def_rpr = list_style.find(qn("a:defRPr"))
                        if def_rpr is not None:
                            replace_properties_fonts(
                                def_rpr,
                                theme_font,
                                preserve_code_fonts,
                                logger,
                            )
        for shape in slide_master.shapes:
            replace_shape_fonts(shape, preserve_code_fonts, logger)
        for j, slide_layout in enumerate(slide_master.slide_layouts):
            logger.log(f"--- Slide Layout {j + 1} ---")
            for shape in slide_layout.shapes:
                replace_shape_fonts(shape, preserve_code_fonts, logger)


def process_notes_master(
    presentation: PresentationType, preserve_code_fonts: bool, logger: Logger
) -> None:
    if presentation.element.find(qn("p:notesMasterIdLst")) is None:
        return
    notes_master = presentation.notes_master
    logger.log("--- Notes Master ---")
    for shape in notes_master.shapes:
        replace_shape_fonts(shape, preserve_code_fonts, logger)


def process_presentation(
    presentation: PresentationType, preserve_code_fonts: bool, logger: Logger
) -> None:
    process_slides(presentation.slides, preserve_code_fonts, logger)
    process_slide_masters(presentation.slide_masters, preserve_code_fonts, logger)
    process_notes_master(presentation, preserve_code_fonts, logger)
